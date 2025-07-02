import os
from pathlib import Path
from datetime import datetime
import pdfplumber
from docx import Document
from pptx import Presentation
import json
from dotenv import load_dotenv

def get_file_metadata(file_path, verbose=True):
    # Convert string path to Path object if needed
    if isinstance(file_path, str):
        file_path = Path(file_path)
    
    if verbose:
        print(f"\n📄 Processing file: {file_path.name}")
        print(f"   Size: {round(os.path.getsize(file_path) / 1024, 2)} KB")
        print(f"   Type: {file_path.suffix.lower()}")
    
    metadata = {
        "filename": file_path.name,
        "extension": file_path.suffix.lower(),
        "size_kb": round(os.path.getsize(file_path) / 1024, 2),
        "created_at": datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
        "modified_at": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
    }

    try:
        ext = file_path.suffix.lower()

        # === PDF ===
        if ext == ".pdf":
            text = ""
            pages_count = 0
            extraction_method = "none"
            
            # Try pdfplumber first (with error handling)
            try:
                with pdfplumber.open(file_path) as pdf:
                    pages_count = len(pdf.pages)
                    text_list = []
                    
                    for page in pdf.pages:
                        try:
                            page_text = page.extract_text()
                            if page_text:
                                text_list.append(page_text.strip())
                        except Exception as page_error:
                            # Skip problematic pages but continue with others
                            continue
                    
                    if text_list:
                        text = "\n".join(text_list)
                        extraction_method = "pdfplumber"
            except Exception as plumber_error:
                # If pdfplumber fails completely, try PyPDF2 as fallback
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        pages_count = len(pdf_reader.pages)
                        text_list = []
                        
                        for page in pdf_reader.pages:
                            try:
                                page_text = page.extract_text()
                                if page_text:
                                    text_list.append(page_text.strip())
                            except Exception:
                                continue
                        
                        if text_list:
                            text = "\n".join(text_list)
                            extraction_method = "PyPDF2"
                except ImportError:
                    # PyPDF2 not available, continue with empty text
                    extraction_method = "failed"
                except Exception:
                    extraction_method = "failed"
            
            word_count = len(text.split()) if text.strip() else 0
            metadata["pages"] = pages_count
            metadata["word_count"] = word_count
            
            if word_count > 30:
                metadata["is_valid"] = True
            else:
                metadata["is_valid"] = False
                if word_count == 0:
                    if extraction_method == "failed":
                        metadata["notes"] = "PDF processing failed - corrupted or encrypted file"
                    else:
                        metadata["notes"] = "Scanned or image-based PDF – no extractable text"
                else:
                    metadata["notes"] = "Low word count for a text-based PDF"

        # === DOCX ===
        elif ext == ".docx":
            doc = Document(file_path)
            texts = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
            
            # Extract text from headers and footers
            for section in doc.sections:
                # Headers
                if section.header:
                    for paragraph in section.header.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text.strip())
                # Footers
                if section.footer:
                    for paragraph in section.footer.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text.strip())
            
            text = "\n".join(texts)
            word_count = len(text.split()) if text.strip() else 0
            metadata["word_count"] = word_count
            metadata["is_valid"] = word_count > 30
            if not metadata["is_valid"]:
                metadata["notes"] = "Too few words in Word document"

        # === PPTX ===
        elif ext == ".pptx":
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            text = "\n".join(texts)
            word_count = len(text.split())
            metadata["slides"] = len(prs.slides)
            metadata["word_count"] = word_count
            metadata["is_valid"] = word_count > 30
            if not metadata["is_valid"]:
                metadata["notes"] = "Too few words in PowerPoint file"

        # === TXT ===
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            word_count = len(text.split())
            metadata["word_count"] = word_count
            metadata["is_valid"] = word_count > 30
            if not metadata["is_valid"]:
                metadata["notes"] = "Text file too short"

        else:
            metadata["is_valid"] = False
            metadata["notes"] = "Unsupported file type"

        # Optional warning for low word count despite large size
        if metadata["word_count"] < 10 and metadata["size_kb"] > 100:
            metadata["notes"] = metadata.get("notes", "") + " File large but almost no text found."

    except Exception as e:
        metadata["is_valid"] = False
        metadata["notes"] = f"Error: {e}"

    return metadata


def process_data_folder(data_folder_path: str, output_folder: str = None) -> list:
    """
    Processes all files in the data folder and collects their metadata.
    Saves results incrementally to a JSON file as each file is processed.
    
    Args:
        data_folder_path: Path to the data folder
        output_folder: Path to save the results (optional)
    
    Returns:
        List of metadata for all files
    """
    data_folder = Path(data_folder_path)
    results = []
    
    if not data_folder.exists() or not data_folder.is_dir():
        print(f"Error: {data_folder_path} is not a valid directory")
        return results
    
    # Set up output file path if output folder is provided
    output_file = None
    if output_folder:
        output_path = Path(output_folder)
        output_path.mkdir(exist_ok=True, parents=True)
        output_file = output_path / f"step1-metadata_{data_folder.name}.json"
        # Initialize with empty list
        with open(output_file, 'w') as f:
            json.dump([], f, indent=4)
    
    for file_path in data_folder.iterdir():
        if file_path.is_file():
            try:
                metadata = get_file_metadata(str(file_path))
                results.append(metadata)
                print(f"Processed: {file_path.name}")
                
                # Save results incrementally after each file
                if output_file:
                    with open(output_file, 'w') as f:
                        json.dump(results, f, indent=4)
                    
            except Exception as e:
                print(f"Error processing {file_path.name}: {str(e)}")
    
    if output_file:
        print(f"Results saved to {output_file}")
    
    return results