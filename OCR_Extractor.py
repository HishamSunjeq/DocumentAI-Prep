import os
import pytesseract
import pdfplumber
from pathlib import Path
from PIL import Image
from docx import Document
from pptx import Presentation
from io import BytesIO
import docx
import pptx
import json
from dotenv import load_dotenv
import shutil
import re
import time

load_dotenv()

# Load cleaning configuration from environment
CLEAN_EXTRACTED_TEXT = os.getenv('CLEAN_EXTRACTED_TEXT', 'true').lower() == 'true'
VERBOSE = os.getenv('VERBOSE_OUTPUT', 'false').lower() == 'true'
DEBUG = os.getenv('DEBUG', 'false').lower() == 'true'

def log_verbose(message, level="info"):
    """Log message only if verbose mode is enabled"""
    if not VERBOSE and level != "error":
        return
        
    prefix = {
        "info": "ℹ️",
        "success": "✅",
        "warning": "⚠️",
        "error": "❌",
        "debug": "🔍",
        "process": "⚙️"
    }.get(level, "ℹ️")
    
    print(f"{prefix} {message}")

# Try to configure tesseract, but handle gracefully if not installed
try:
    tesseract_path = shutil.which("tesseract")
    
    # If not found in PATH, check common Windows installation locations
    if not tesseract_path:
        common_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Users\{}\AppData\Local\Tesseract-OCR\tesseract.exe".format(os.getenv('USERNAME', ''))
        ]
        
        for path in common_paths:
            if os.path.exists(path):
                tesseract_path = path
                break
    
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
        log_verbose(f"Tesseract found at: {tesseract_path}", "success")
    else:
        log_verbose("Tesseract not found in PATH or common locations.", "warning")
        log_verbose("To install Tesseract on Windows:", "info")
        log_verbose("1. Download from: https://github.com/UB-Mannheim/tesseract/wiki", "info")
        log_verbose("2. Install and add to PATH", "info")
        log_verbose("3. Or use: winget install tesseract-ocr.tesseract", "info")
except Exception as e:
    log_verbose(f"Error configuring Tesseract: {e}", "warning")
    log_verbose("OCR functionality will be limited.", "warning")

def is_garbage_line(line):
    """Check if a line appears to be garbage/corrupted text from OCR errors."""
    # Line is very short and not meaningful
    if len(line.strip()) < 8:
        return True

    # Contains too many special characters or digits
    if sum(c.isalnum() for c in line) / (len(line) + 1e-6) < 0.4:
        return True

    # Contains too many random capitalized words (common in OCR errors)
    if len(re.findall(r'\b[A-Z]{2,}\b', line)) > 3:
        return True

    # Looks like a garbled line with mixed letters/digits/symbols
    if re.search(r'[^\w\s]{2,}', line):
        return True

    return False

def clean_text(text):
    """Clean and normalize extracted text by removing common artifacts and formatting issues."""
    if not text:
        return ""
    
    # Remove page/slide numbers
    text = re.sub(r'Page\s+\d+\s+of\s+\d+', '', text, flags=re.IGNORECASE)
    text = re.sub(r'Slide\s+\d+', '', text, flags=re.IGNORECASE)
    
    # Remove company watermarks (adjust pattern as needed)
    text = re.sub(r'BAYANAT\s+\(?\d{4}\)?', '', text, flags=re.IGNORECASE)
    
    # Normalize whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)  # Replace 3+ newlines with 2
    text = re.sub(r'\s{2,}', ' ', text)     # Replace multiple spaces with single space
    
    # Remove bullet points and numbering at start of lines
    text = re.sub(r'^[\-\*\d\.\)]\s+', '', text, flags=re.MULTILINE)
    
    # Remove non-ASCII characters (keeping basic punctuation)
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)
    
    # Remove extra whitespace from line breaks
    text = re.sub(r'\s*\n\s*', '\n', text)
    
    # Remove email addresses
    text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[email_redacted]', text)
    
     # Remove Table of Contents lines (based on dots and page numbers)
    text = re.sub(r'\.{3,}\s*\d{1,3}', '', text)
    
    # Remove common confidential keywords
    text = re.sub(r'Confidential!?|Internal Use Only|Draft Copy', '', text, flags=re.IGNORECASE)
    
     # Remove "Table of Contents" section titles
    text = re.sub(r'Table of Contents', '', text, flags=re.IGNORECASE)
    
    # Filter out garbage lines using line-by-line analysis
    lines = text.split('\n')
    clean_lines = []
    
    for line in lines:
        if not is_garbage_line(line):
            clean_lines.append(line)
        # Optional: Log removed garbage lines for debugging
        # else:
        #     print(f"🗑️ Removed garbage line: {line[:50]}...")
    
    text = '\n'.join(clean_lines)
    
    # Collapse multiple spaces and newlines
    text = re.sub(r'\n{2,}', '\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    
    # Remove leading/trailing whitespace
    return text.strip()

def ocr_image(image: Image.Image):
    if image is None:
        log_verbose("Skipped OCR: Invalid or None image input", "warning")
        return ""
    
    try:
        # Check if tesseract is available
        if not pytesseract.pytesseract.tesseract_cmd:
            log_verbose("Skipped OCR: Tesseract not available", "warning")
            return ""
        
        # Validate image dimensions
        if not hasattr(image, 'size') or len(image.size) != 2:
            log_verbose("Skipped OCR: Invalid image dimensions", "warning")
            return ""
        
        width, height = image.size
        if width < 10 or height < 10:
            log_verbose("Skipped OCR: Image too small for text extraction", "warning")
            return ""
        
        # Ensure image is in a compatible mode for OCR
        if image.mode not in ['RGB', 'L', 'RGBA']:
            try:
                image = image.convert('RGB')
            except Exception as e:
                log_verbose(f"Failed to convert image mode: {e}", "warning")
                return ""
        
        # Perform OCR
        result = pytesseract.image_to_string(image, config='--psm 6')
        return result.strip()
        
    except Exception as e:
        log_verbose(f"OCR failed: {e}", "warning")
        return ""

def extract_pdf_ocr(path):
    text_output = []
    
    try:
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                page_text = ""
                
                # Try direct text extraction first
                try:
                    direct_text = page.extract_text()
                    if direct_text and len(direct_text.strip()) > 10:
                        page_text = direct_text
                except Exception as e:
                    log_verbose(f"Failed direct text extraction on page {page_num + 1}: {e}", "warning")
                
                # Try table extraction if direct text is insufficient
                if not page_text or len(page_text.strip()) < 10:
                    try:
                        tables = page.extract_tables()
                        if tables:
                            table_texts = []
                            for table in tables:
                                if table:
                                    table_rows = []
                                    for row in table:
                                        if row:
                                            clean_row = [str(cell).strip() if cell else "" for cell in row]
                                            if any(clean_row):  # Only add rows with content
                                                table_rows.append(" | ".join(clean_row))
                                    if table_rows:
                                        table_texts.append("\n".join(table_rows))
                            if table_texts:
                                page_text = "\n\n".join(table_texts)
                    except Exception as e:
                        log_verbose(f"Failed table extraction on page {page_num + 1}: {e}", "warning")
                
                # OCR fallback (only if tesseract is available and no text was found)
                if (not page_text or len(page_text.strip()) < 10) and pytesseract.pytesseract.tesseract_cmd:
                    try:
                        log_verbose(f"Using OCR for page {page_num + 1} (low/no text content detected)", "process")
                        image = page.to_image(resolution=300).original
                        if image is not None:
                            ocr_text = ocr_image(image)
                            if ocr_text and len(ocr_text.strip()) > len(page_text.strip()):
                                page_text = ocr_text
                    except Exception as e:
                        log_verbose(f"Failed to convert PDF page {page_num + 1} to image: {e}", "warning")
                
                if page_text:
                    text_output.append(page_text)
                    
    except Exception as e:
        log_verbose(f"Failed to process PDF: {e}", "error")
        return ""
    
    return "\n\n".join(text_output).strip()

def extract_docx_ocr(path):
    doc = Document(path)
    all_text = []
    
    # Extract text from paragraphs
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if paragraphs:
        all_text.extend(paragraphs)
    
    # Extract text from tables
    try:
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    table_text.append(" | ".join(row_text))
            if table_text:
                all_text.append("\n".join(table_text))
    except Exception as e:
        log_verbose(f"Failed to extract table content: {e}", "warning")
    
    # Extract text from headers
    try:
        for section in doc.sections:
            if section.header:
                header_text = []
                for paragraph in section.header.paragraphs:
                    if paragraph.text.strip():
                        header_text.append(paragraph.text.strip())
                if header_text:
                    all_text.extend(header_text)
    except Exception as e:
        log_verbose(f"Failed to extract header content: {e}", "warning")
    
    # Extract text from footers
    try:
        for section in doc.sections:
            if section.footer:
                footer_text = []
                for paragraph in section.footer.paragraphs:
                    if paragraph.text.strip():
                        footer_text.append(paragraph.text.strip())
                if footer_text:
                    all_text.extend(footer_text)
    except Exception as e:
        log_verbose(f"Failed to extract footer content: {e}", "warning")
    
    # Extract text from text boxes and shapes
    try:
        # This requires accessing the document's XML structure
        from docx.oxml.ns import qn
        from docx.oxml import parse_xml
        
        # Get all text from text boxes and shapes in the document
        for element in doc.element.iter():
            if element.tag.endswith('}t'):  # Text elements
                if element.text and element.text.strip():
                    # Check if this text is not already captured in paragraphs/tables
                    text = element.text.strip()
                    if text not in " ".join(all_text):
                        all_text.append(text)
    except Exception as e:
        log_verbose(f"Failed to extract text from shapes/textboxes: {e}", "warning")
    
    # Combine all extracted text
    full_text = "\n".join(all_text)
    
    # OCR embedded images (e.g., scanned documents in .docx)
    try:
        if hasattr(doc.part, '_rels') and doc.part._rels:
            for rel_id in doc.part._rels:
                rel = doc.part._rels[rel_id]
                if hasattr(rel, 'target_ref') and "image" in rel.target_ref:
                    try:
                        if not hasattr(rel, 'target_part') or rel.target_part is None:
                            continue
                        
                        image_data = rel.target_part.blob
                        if not image_data or len(image_data) == 0:
                            continue
                        
                        # Try to open the image
                        image = Image.open(BytesIO(image_data))
                        
                        # Validate image format and size
                        if image.format == 'WMF':
                            # WMF files are vector format, not suitable for OCR - skip silently
                            continue
                        elif image.format not in ['JPEG', 'PNG', 'BMP', 'TIFF', 'GIF']:
                            log_verbose(f"Skipping unsupported image format: {image.format}", "warning")
                            continue
                        
                        if image.size[0] < 50 or image.size[1] < 50:  # Skip very small images
                            continue
                        
                        # Convert to RGB if necessary (for OCR compatibility)
                        if image.mode not in ['RGB', 'L']:
                            image = image.convert('RGB')
                        
                        ocr_text = ocr_image(image)
                        if ocr_text and ocr_text.strip():
                            full_text += "\n" + ocr_text.strip()
                            
                    except Exception as e:
                        log_verbose(f"Failed to extract image from DOCX: {e}", "warning")
                        continue
    except Exception as e:
        log_verbose(f"Failed to process DOCX images: {e}", "warning")
    
    return full_text.strip()

def extract_pptx_ocr(path):
    try:
        prs = Presentation(path)
        all_texts = []

        for slide_num, slide in enumerate(prs.slides):
            slide_texts = []
            
            # Extract text from all shapes with comprehensive approach
            for shape in slide.shapes:
                try:
                    # Approach 1: Direct text extraction (works for most text shapes)
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                        continue
                    
                    # Approach 2: Text frame extraction (backup method)
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        try:
                            if hasattr(shape.text_frame, "text") and shape.text_frame.text.strip():
                                slide_texts.append(shape.text_frame.text.strip())
                                continue
                            
                            # Extract from individual paragraphs in text frame
                            paragraph_texts = []
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text and paragraph.text.strip():
                                    paragraph_texts.append(paragraph.text.strip())
                            if paragraph_texts:
                                slide_texts.append("\n".join(paragraph_texts))
                                continue
                        except:
                            pass
                    
                    # Approach 3: Handle specific shape types
                    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                        try:
                            if hasattr(shape, "table") and shape.table:
                                table_text = []
                                for row in shape.table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        if cell.text and cell.text.strip():
                                            row_text.append(cell.text.strip())
                                    if row_text:
                                        table_text.append(" | ".join(row_text))
                                if table_text:
                                    slide_texts.append("\n".join(table_text))
                                    continue
                        except Exception as table_e:
                            log_verbose(f"Failed to extract table from slide {slide_num + 1}: {table_e}", "warning")
                    
                    elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                        try:
                            if hasattr(shape, 'shapes'):  # Group shape
                                group_texts = []
                                for grouped_shape in shape.shapes:
                                    # Try multiple methods for grouped shapes
                                    group_text = None
                                    
                                    # Method 1: Direct text
                                    if hasattr(grouped_shape, "text") and grouped_shape.text and grouped_shape.text.strip():
                                        group_text = grouped_shape.text.strip()
                                    # Method 2: Text frame
                                    elif hasattr(grouped_shape, "text_frame") and grouped_shape.text_frame:
                                        if hasattr(grouped_shape.text_frame, "text") and grouped_shape.text_frame.text.strip():
                                            group_text = grouped_shape.text_frame.text.strip()
                                        else:
                                            # Method 3: Individual paragraphs in text frame
                                            para_texts = []
                                            try:
                                                for para in grouped_shape.text_frame.paragraphs:
                                                    if para.text and para.text.strip():
                                                        para_texts.append(para.text.strip())
                                                if para_texts:
                                                    group_text = "\n".join(para_texts)
                                            except:
                                                pass
                                    
                                    if group_text:
                                        group_texts.append(group_text)
                                
                                if group_texts:
                                    slide_texts.append("\n".join(group_texts))
                                    continue
                        except Exception as group_e:
                            log_verbose(f"Failed to extract text from grouped shapes in slide {slide_num + 1}: {group_e}", "warning")
                    
                    elif shape.shape_type == 13:  # Picture/Image - OCR
                        try:
                            # Check if shape actually has an image
                            if not hasattr(shape, 'image') or shape.image is None:
                                continue
                            
                            image_stream = shape.image.blob
                            if not image_stream or len(image_stream) == 0:
                                continue
                            
                            # Try to open the image
                            image = Image.open(BytesIO(image_stream))
                            
                            # Validate image format and size
                            if image.format == 'WMF':
                                # WMF files are vector format, not suitable for OCR - skip silently
                                continue
                            elif image.format not in ['JPEG', 'PNG', 'BMP', 'TIFF', 'GIF']:
                                log_verbose(f"Skipping unsupported image format: {image.format}", "warning")
                                continue
                            
                            if image.size[0] < 50 or image.size[1] < 50:  # Skip very small images
                                continue
                            
                            # Convert to RGB if necessary (for OCR compatibility)
                            if image.mode not in ['RGB', 'L']:
                                image = image.convert('RGB')
                            
                            ocr_text = ocr_image(image)
                            if ocr_text and ocr_text.strip():
                                slide_texts.append(f"[Image OCR]: {ocr_text.strip()}")
                                
                        except AttributeError:
                            # Shape doesn't have image attribute
                            continue
                        except Exception as e:
                            log_verbose(f"Failed to extract image from slide {slide_num + 1}: {e}", "warning")
                            continue
                
                except AttributeError as attr_e:
                    # Shape doesn't have expected attributes - this is normal, skip silently
                    continue
                except Exception as e:
                    # Only log unexpected errors
                    if "shape does not contain a table" not in str(e):
                        log_verbose(f"Failed to process shape in slide {slide_num + 1}: {e}", "warning")
                    continue
            
            # Add slide content if any text was found
            if slide_texts:
                slide_content = f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_texts)
                all_texts.append(slide_content)
        
        return "\n\n".join(all_texts).strip()
        
    except Exception as e:
        log_verbose(f"Failed to process PowerPoint file: {e}", "error")
        return ""

def extract_txt(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except:
        return ""  # fallback to OCR only if needed

def extract_text_from_any_file(file_path):
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return extract_pdf_ocr(file_path)
    elif ext == ".docx":
        return extract_docx_ocr(file_path)
    elif ext == ".pptx":
        return extract_pptx_ocr(file_path)
    elif ext == ".txt":
        return extract_txt(file_path)
    else:
        return None

def save_extracted_text(file_path, extracted_text, output_folder="./output/ocr_extracted"):
    """Save extracted text to a file in the output folder."""
    if not extracted_text or not extracted_text.strip():
        return None
    
    # Apply cleaning only if enabled in environment settings
    if CLEAN_EXTRACTED_TEXT:
        final_text = clean_text(extracted_text)
        text_status = "cleaned"
    else:
        final_text = extracted_text
        text_status = "raw"
    
    if not final_text or not final_text.strip():
        log_verbose(f"No text remaining after {'cleaning' if CLEAN_EXTRACTED_TEXT else 'processing'} for {Path(file_path).name}", "warning")
        return None
    
    # Create output directory if it doesn't exist
    output_path = Path(output_folder)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Create output filename (same name but with .txt extension)
    input_file = Path(file_path)
    output_filename = f"{input_file.stem}_extracted.txt"
    output_file = output_path / output_filename
    
    # Save the text
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(final_text)
        log_verbose(f"Saved {text_status} text: {output_file}", "success")
        return str(output_file)
    except Exception as e:
        log_verbose(f"Error saving extracted text for {input_file.name}: {e}", "error")
        return None

def process_file_with_ocr(file_path, output_folder="./output/ocr_extracted"):
    """Extract text from a single file and save it."""
    file_path = Path(file_path)
    
    log_verbose(f"Processing: {file_path.name}", "process")
    
    # Extract text using OCR
    extracted_text = extract_text_from_any_file(file_path)
    
    if extracted_text and extracted_text.strip():
        # Save the extracted text (cleaning happens conditionally inside save_extracted_text)
        saved_file = save_extracted_text(file_path, extracted_text, output_folder)
        
        if saved_file:
            # Use final text (cleaned or raw) for statistics based on settings
            final_text = clean_text(extracted_text) if CLEAN_EXTRACTED_TEXT else extracted_text
            word_count = len(final_text.split()) if final_text else 0
            char_count = len(final_text) if final_text else 0
            
            log_verbose(f"Extracted {char_count} characters ({word_count} words) from {file_path.name} ({'cleaned' if CLEAN_EXTRACTED_TEXT else 'raw'} text)", "info")

            # Create preview from final text
            preview = final_text[:200] + "..." if len(final_text) > 200 else final_text
            
            return {
                "input_file": str(file_path),
                "output_file": saved_file,
                "word_count": word_count,
                "char_count": char_count,
                "success": True,
                "preview": preview,
                "text_cleaned": CLEAN_EXTRACTED_TEXT
            }
        else:
            log_verbose(f"Failed to save extracted text from {file_path.name}", "error")
            return {
                "input_file": str(file_path),
                "output_file": None,
                "word_count": 0,
                "char_count": 0,
                "success": False,
                "preview": "",
                "text_cleaned": CLEAN_EXTRACTED_TEXT
            }
    else:
        log_verbose(f"No text extracted from {file_path.name}", "error")
        return {
            "input_file": str(file_path),
            "output_file": None,
            "word_count": 0,
            "char_count": 0,
            "success": False,
            "preview": "",
            "text_cleaned": CLEAN_EXTRACTED_TEXT
        }

def process_folder_with_ocr(data_folder, output_folder):
    """Process all files in a folder using OCR and save results."""
    data_path = Path(data_folder)
    results = []
    
    if not data_path.exists():
        log_verbose(f"Error: Data folder '{data_folder}' does not exist", "error")
        return results
    
    log_verbose(f"Processing files from: {data_path}")
    log_verbose(f"Saving extracted text to: {output_folder}")
    log_verbose(f"Text cleaning: {'Enabled' if CLEAN_EXTRACTED_TEXT else 'Disabled'}")
    
    # Get all supported files (excluding temporary Microsoft Office files)
    supported_extensions = {'.pdf', '.docx', '.pptx', '.txt'}
    files_to_process = [f for f in data_path.iterdir() 
                       if f.is_file() and f.suffix.lower() in supported_extensions
                       and not f.name.startswith('~$')]  # Skip temporary Office files
    
    if not files_to_process:
        log_verbose("No supported files found in the data folder", "warning")
        return results
    
    # Process each file
    for file_path in files_to_process:
        try:
            result = process_file_with_ocr(file_path, output_folder)
            results.append(result)
        except Exception as e:
            log_verbose(f"Error processing {file_path.name}: {e}", "error")
            results.append({
                "input_file": str(file_path),
                "output_file": None,
                "word_count": 0,
                "char_count": 0,
                "success": False,
                "error": str(e),
                "text_cleaned": CLEAN_EXTRACTED_TEXT
            })
    
    # Save processing summary
    summary_file = Path(output_folder) / "ocr_processing_summary.json"
    try:
        with open(summary_file, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=4)
        log_verbose(f"Processing summary saved to: {summary_file}", "success")
    except Exception as e:
        log_verbose(f"Error saving summary: {e}", "error")
    
    # Print summary
    successful = sum(1 for r in results if r['success'])
    total_words = sum(r['word_count'] for r in results if r['success'])
    
    log_verbose(f"\nOCR Processing Complete:")
    log_verbose(f"- Files processed: {len(results)}")
    log_verbose(f"- Successful extractions: {successful}")
    log_verbose(f"- Total words extracted: {total_words}")
    log_verbose(f"- Text cleaning: {'Enabled' if CLEAN_EXTRACTED_TEXT else 'Disabled'}")
    
    return results


if __name__ == "__main__":
    # Process the data folder using OCR
    data_folder = "./data"
    output_folder = "./output/ocr_output"
    process_folder_with_ocr(data_folder, output_folder)